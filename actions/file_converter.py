import io
import json
import logging
import os
import re
import subprocess
import tempfile
import traceback
from pathlib import Path

try:
    from PIL import Image
    import pytesseract
except ImportError:
    pytesseract = None

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from openpyxl import load_workbook
    from openpyxl.utils import get_column_letter
except ImportError:
    load_workbook = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import markdown
except ImportError:
    markdown = None

try:
    import img2pdf
except ImportError:
    img2pdf = None

try:
    from pdf2image import convert_from_path
except ImportError:
    convert_from_path = None

try:
    import weasyprint
except ImportError:
    weasyprint = None

try:
    from pydub import AudioSegment
except ImportError:
    AudioSegment = None

logger = logging.getLogger("file_converter")

_TEMP_DIR = Path(tempfile.gettempdir()) / "jarvis_converter"
_TEMP_DIR.mkdir(parents=True, exist_ok=True)

SUPPORTED_IMAGE_READ = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".tif", ".webp", ".ico"}
SUPPORTED_IMAGE_WRITE = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".tif", ".webp", ".ico"}

_EXT_MAP = {
    "png": ".png", "jpg": ".jpg", "jpeg": ".jpg", "gif": ".gif",
    "bmp": ".bmp", "tiff": ".tiff", "tif": ".tiff", "webp": ".webp",
    "ico": ".ico", "pdf": ".pdf", "docx": ".docx", "doc": ".docx",
    "xlsx": ".xlsx", "xls": ".xlsx", "pptx": ".pptx", "ppt": ".pptx",
    "md": ".md", "html": ".html", "htm": ".html", "csv": ".csv",
    "json": ".json", "xml": ".xml", "txt": ".txt",
    "mp3": ".mp3", "wav": ".wav", "ogg": ".ogg", "flac": ".flac",
    "m4a": ".m4a", "aac": ".aac",
    "mp4": ".mp4", "avi": ".avi", "mkv": ".mkv", "mov": ".mov",
    "webm": ".webm", "gif": ".gif",
}


def _guess_ext(fmt: str) -> str:
    return _EXT_MAP.get(fmt.lower().strip(" ."), "." + fmt.lower().strip(" ."))


def _output_path(source: str, target_ext: str) -> str:
    src = Path(source)
    name = src.stem
    out = _TEMP_DIR / f"{name}_{target_ext.lstrip('.')}{target_ext}"
    return str(out)


def _open_docx_text(path: str) -> str:
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs)


def _open_pptx_text(path: str) -> str:
    prs = Presentation(path)
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    lines.append(para.text)
    return "\n".join(lines)


def _open_xlsx_text(path: str) -> str:
    wb = load_workbook(path, read_only=True, data_only=True)
    lines = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            lines.append(",".join(str(c) if c is not None else "" for c in row))
    wb.close()
    return "\n".join(lines)


def convert_image(source: str, target_ext: str, **kwargs) -> str:
    img = Image.open(source)
    out = _output_path(source, target_ext)
    if target_ext.lower() in (".jpg", ".jpeg"):
        img = img.convert("RGB")
    img.save(out)
    return out


def convert_image_to_pdf(source: str, **kwargs) -> str:
    if img2pdf is None:
        raise RuntimeError("img2pdf not installed")
    img = Image.open(source)
    if img.mode != "RGB":
        img = img.convert("RGB")
    out = _output_path(source, ".pdf")
    with open(out, "wb") as f:
        f.write(img2pdf.convert(source))
    return out


def convert_pdf_to_image(source: str, target_ext: str = ".png", **kwargs) -> str:
    if convert_from_path is None:
        raise RuntimeError("pdf2image not installed")
    images = convert_from_path(source, fmt=target_ext.lstrip(".").lower())
    out_dir = _TEMP_DIR / Path(source).stem
    out_dir.mkdir(parents=True, exist_ok=True)
    paths = []
    for i, img in enumerate(images):
        p = str(out_dir / f"page_{i+1:03d}{target_ext}")
        img.save(p)
        paths.append(p)
    return paths[0] if len(paths) == 1 else "\n".join(paths)


def convert_pdf_to_docx(source: str, **kwargs) -> str:
    if fitz is None:
        raise RuntimeError("PyMuPDF not installed")
    from pdf2docx import Converter
    out = _output_path(source, ".docx")
    cv = Converter(source)
    cv.convert(out, start=0, end=None)
    cv.close()
    return out


def convert_pdf_to_xlsx(source: str, **kwargs) -> str:
    if fitz is None:
        raise RuntimeError("PyMuPDF not installed")
    import pdfplumber
    out = _output_path(source, ".xlsx")
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "PDF Data"
    with pdfplumber.open(source) as pdf:
        row = 1
        for page in pdf.pages:
            tables = page.extract_tables()
            for table in tables:
                for r in table:
                    for col_idx, val in enumerate(r, 1):
                        ws.cell(row=row, column=col_idx, value=val)
                    row += 1
                row += 1
            text = page.extract_text()
            if text:
                for line in text.split("\n"):
                    ws.cell(row=row, column=1, value=line)
                    row += 1
    wb.save(out)
    wb.close()
    return out


def convert_docx_to_pdf(source: str, **kwargs) -> str:
    text = _open_docx_text(source)
    return _html_to_pdf(text, source, ".pdf")


def convert_xlsx_to_pdf(source: str, **kwargs) -> str:
    text = _open_xlsx_text(source)
    html = "<html><body><pre>" + text + "</pre></body></html>"
    return _html_to_pdf_raw(html, source, ".pdf")


def convert_pptx_to_pdf(source: str, **kwargs) -> str:
    text = _open_pptx_text(source)
    return _html_to_pdf(text, source, ".pdf")


def convert_md_to_html(source: str, **kwargs) -> str:
    raw = Path(source).read_text(encoding="utf-8")
    if markdown is None:
        raise RuntimeError("markdown library not installed")
    html = markdown.markdown(raw, extensions=["extra", "codehilite"])
    out = _output_path(source, ".html")
    full = f"<!DOCTYPE html><html><head><meta charset='utf-8'></head><body>{html}</body></html>"
    Path(out).write_text(full, encoding="utf-8")
    return out


def convert_html_to_pdf(source: str, **kwargs) -> str:
    if weasyprint is None:
        raise RuntimeError("weasyprint not installed")
    out = _output_path(source, ".pdf")
    weasyprint.HTML(filename=source).write_pdf(out)
    return out


def _html_to_pdf(text: str, source: str, target_ext: str) -> str:
    html = f"<!DOCTYPE html><html><head><meta charset='utf-8'></head><body><pre>{text}</pre></body></html>"
    return _html_to_pdf_raw(html, source, target_ext)


def _html_to_pdf_raw(html: str, source: str, target_ext: str) -> str:
    if weasyprint is None:
        raise RuntimeError("weasyprint not installed")
    out = _output_path(source, target_ext)
    weasyprint.HTML(string=html).write_pdf(out)
    return out


def convert_text_to_csv(source: str, **kwargs) -> str:
    raw = Path(source).read_text(encoding="utf-8")
    import csv
    out = _output_path(source, ".csv")
    lines = [l for l in raw.split("\n") if l.strip()]
    with open(out, "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        for line in lines:
            writer.writerow(line.split("\t" if "\t" in line else ","))
    return out


def convert_csv_to_json(source: str, **kwargs) -> str:
    import csv
    out = _output_path(source, ".json")
    with open(source, newline="", encoding="utf-8") as f:
        reader = csv.DictReader(f)
        data = list(reader)
    Path(out).write_text(json.dumps(data, indent=2), encoding="utf-8")
    return out


def convert_json_to_csv(source: str, **kwargs) -> str:
    import csv
    out = _output_path(source, ".csv")
    with open(source, encoding="utf-8") as f:
        data = json.load(f)
    if not data:
        raise RuntimeError("Empty JSON array")
    import openpyxl
    with open(out, "w", newline="", encoding="utf-8") as f:
        writer = csv.DictWriter(f, fieldnames=list(data[0].keys()))
        writer.writeheader()
        writer.writerows(data)
    return out


def convert_xml_to_json(source: str, **kwargs) -> str:
    import xml.etree.ElementTree as ET
    out = _output_path(source, ".json")

    def elem_to_dict(e):
        d = {}
        for child in e:
            val = elem_to_dict(child) if len(child) else child.text or ""
            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
            if tag in d:
                if not isinstance(d[tag], list):
                    d[tag] = [d[tag]]
                d[tag].append(val)
            else:
                d[tag] = val
        return d

    tree = ET.parse(source)
    root = tree.getroot()
    tag = root.tag.split("}")[-1] if "}" in root.tag else root.tag
    result = {tag: elem_to_dict(root)}
    Path(out).write_text(json.dumps(result, indent=2), encoding="utf-8")
    return out


def convert_video_to_gif(source: str, **kwargs) -> str:
    out = _output_path(source, ".gif")
    fps = kwargs.get("fps", 10)
    width = kwargs.get("width", 320)
    subprocess.run([
        "ffmpeg", "-y", "-i", source,
        "-vf", f"fps={fps},scale={width}:-1:flags=lanczos",
        "-c:v", "gif", out
    ], capture_output=True, check=True)
    return out


def convert_video(source: str, target_ext: str, **kwargs) -> str:
    out = _output_path(source, target_ext)
    subprocess.run([
        "ffmpeg", "-y", "-i", source, out
    ], capture_output=True, check=True)
    return out


def convert_audio(source: str, target_ext: str, **kwargs) -> str:
    if AudioSegment is None:
        raise RuntimeError("pydub not installed")
    out = _output_path(source, target_ext)
    audio = AudioSegment.from_file(source)
    audio.export(out, format=target_ext.lstrip("."))
    return out


def convert_image_ocr(source: str, **kwargs) -> str:
    if pytesseract is None:
        raise RuntimeError("pytesseract not installed. Install with: pip install pytesseract")
    img = Image.open(source)
    text = pytesseract.image_to_string(img)
    out = _output_path(source, ".txt")
    Path(out).write_text(text, encoding="utf-8")
    return out


_CONVERTERS = {
    # Image → *
    (".png", ".pdf"): convert_image_to_pdf,
    (".jpg", ".pdf"): convert_image_to_pdf,
    (".jpeg", ".pdf"): convert_image_to_pdf,
    (".gif", ".pdf"): convert_image_to_pdf,
    (".bmp", ".pdf"): convert_image_to_pdf,
    (".tiff", ".pdf"): convert_image_to_pdf,
    (".webp", ".pdf"): convert_image_to_pdf,
    (".jpg", ".png"): convert_image, (".png", ".jpg"): convert_image,
    (".png", ".webp"): convert_image, (".webp", ".png"): convert_image,
    (".png", ".gif"): convert_image, (".gif", ".png"): convert_image,
    (".png", ".bmp"): convert_image, (".png", ".ico"): convert_image,
    (".png", ".tiff"): convert_image, (".tiff", ".png"): convert_image,
    # PDF → *
    (".pdf", ".png"): convert_pdf_to_image,
    (".pdf", ".jpg"): convert_pdf_to_image,
    (".pdf", ".docx"): convert_pdf_to_docx,
    (".pdf", ".xlsx"): convert_pdf_to_xlsx,
    # Docx → *
    (".docx", ".pdf"): convert_docx_to_pdf,
    (".doc", ".pdf"): convert_docx_to_pdf,
    (".docx", ".txt"): lambda s, **kw: str(Path(_output_path(s, ".txt")).write_text(_open_docx_text(s), encoding="utf-8")) or _output_path(s, ".txt"),
    # Xlsx → *
    (".xlsx", ".pdf"): convert_xlsx_to_pdf,
    (".xls", ".pdf"): convert_xlsx_to_pdf,
    # Pptx → *
    (".pptx", ".pdf"): convert_pptx_to_pdf,
    (".ppt", ".pdf"): convert_pptx_to_pdf,
    # Markdown → *
    (".md", ".html"): convert_md_to_html,
    # HTML → *
    (".html", ".pdf"): convert_html_to_pdf,
    (".htm", ".pdf"): convert_html_to_pdf,
    # CSV → *
    (".csv", ".json"): convert_csv_to_json,
    # JSON → *
    (".json", ".csv"): convert_json_to_csv,
    # XML → *
    (".xml", ".json"): convert_xml_to_json,
    # Text → *
    (".txt", ".csv"): convert_text_to_csv,
    # Video → GIF
    (".mp4", ".gif"): convert_video_to_gif,
    (".avi", ".gif"): convert_video_to_gif,
    (".mkv", ".gif"): convert_video_to_gif,
    (".mov", ".gif"): convert_video_to_gif,
    (".webm", ".gif"): convert_video_to_gif,
    # OCR
    (".png", ".txt"): convert_image_ocr,
    (".jpg", ".txt"): convert_image_ocr,
    (".jpeg", ".txt"): convert_image_ocr,
    (".bmp", ".txt"): convert_image_ocr,
    (".tiff", ".txt"): convert_image_ocr,
    (".webp", ".txt"): convert_image_ocr,
}

_VIDEO_EXTENSIONS = {".mp4", ".avi", ".mkv", ".mov", ".webm"}
_AUDIO_EXTENSIONS = {".mp3", ".wav", ".ogg", ".flac", ".m4a", ".aac"}


def convert_file(parameters: dict = None, player=None) -> str:
    params = parameters or {}
    source = params.get("source_path", "")
    target_fmt = params.get("target_format", "").lower().strip(" .")
    mode = params.get("mode", "auto")

    if not source:
        return "❌ No source file specified. Use `source_path` parameter with the full file path."

    if not os.path.isfile(source):
        return f"❌ Source file not found: {source}"

    src_ext = Path(source).suffix.lower()
    target_ext = _guess_ext(target_fmt) if target_fmt else ".pdf"

    if not target_ext:
        return f"❌ Unknown target format: {target_fmt}"

    if src_ext == target_ext:
        return f"Source and target are the same format ({src_ext}). Nothing to do."

    key = (src_ext, target_ext)

    if key in _CONVERTERS:
        try:
            func = _CONVERTERS[key]
            try:
                result = func(source, target_ext, player=player)
            except TypeError:
                result = func(source, player=player)
            msg = f"✅ Converted {Path(source).name} → {result}"
            if player:
                try:
                    player.speak(msg)
                except Exception:
                    pass
            return msg
        except Exception as e:
            logger.error("Conversion failed: %s\n%s", e, traceback.format_exc())
            return f"❌ Conversion failed: {e}"

    # Try generic converters
    try:
        # Image → image
        if src_ext in SUPPORTED_IMAGE_READ and target_ext in SUPPORTED_IMAGE_WRITE:
            out = convert_image(source, target_ext)
            return f"✅ Converted {Path(source).name} → {out}"

        # Video → video
        if src_ext in _VIDEO_EXTENSIONS and target_ext in _VIDEO_EXTENSIONS:
            out = convert_video(source, target_ext)
            return f"✅ Converted {Path(source).name} → {out}"

        # Audio → audio
        if src_ext in _AUDIO_EXTENSIONS and target_ext in _AUDIO_EXTENSIONS:
            out = convert_audio(source, target_ext)
            return f"✅ Converted {Path(source).name} → {out}"

        # Image OCR (any image → text)
        if src_ext in SUPPORTED_IMAGE_READ and target_ext == ".txt":
            out = convert_image_ocr(source)
            return f"✅ OCR extracted text → {out}"

    except Exception as e:
        logger.error("Conversion failed: %s\n%s", e, traceback.format_exc())
        return f"❌ Conversion failed: {e}"

    return f"❌ Unsupported conversion: {src_ext} → {target_ext}. Available: image format swap, PDF↔image, PDF→docx, docx→pdf, xlsx→pdf, pptx→pdf, md→html, html→pdf, csv→json, json→csv, xml→json, txt→csv, video→gif, OCR (image→text), and more."
