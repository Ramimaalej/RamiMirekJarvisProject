"""
Deep comprehensive tests for every action module — run with:
    pytest tests/test_deep.py -x -v --tb=short  (stop on first failure)
    pytest tests/test_deep.py -v --tb=long       (full traceback)
"""

import csv
import json
import os
import re
import sys
import tempfile
from pathlib import Path

import pytest

PROJECT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(PROJECT))


def _import(name, path):
    import importlib.util
    spec = importlib.util.spec_from_file_location(name, path)
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)
    return mod


# ═══════════════════════════════════════════════════════════════════════
# FIXTURES
# ═══════════════════════════════════════════════════════════════════════


@pytest.fixture
def tmp_workspace():
    path = Path(tempfile.mkdtemp(prefix="test_converter_"))
    yield path
    import shutil
    shutil.rmtree(path, ignore_errors=True)


@pytest.fixture
def mock_player():
    class MockPlayer:
        def __init__(self):
            self.last_url = None
        def open_todo_panel(self, url):
            self.last_url = url
    return MockPlayer()


@pytest.fixture
def isolated_tasks(monkeypatch):
    tmp = Path(tempfile.mktemp(suffix=".json"))
    monkeypatch.setattr("actions.task_manager.TASKS_PATH", tmp)
    yield tmp
    tmp.unlink(missing_ok=True)


@pytest.fixture
def with_sample_image(tmp_workspace):
    path = tmp_workspace / "test_img.png"
    try:
        from PIL import Image
        Image.new("RGB", (50, 50), color=(255, 0, 0)).save(str(path))
    except ImportError:
        pytest.skip("Pillow not installed")
    return path


@pytest.fixture
def with_sample_csv(tmp_workspace):
    path = tmp_workspace / "test_data.csv"
    with open(str(path), "w", newline="") as f:
        w = csv.writer(f)
        w.writerow(["name", "age", "city"])
        w.writerow(["Alice", "30", "Paris"])
        w.writerow(["Bob", "25", "London"])
    return path


@pytest.fixture
def with_sample_md(tmp_workspace):
    path = tmp_workspace / "test.md"
    path.write_text("# Hello\n\nThis is **markdown**.\n\n- item 1\n- item 2")
    return path


@pytest.fixture
def with_sample_txt(tmp_workspace):
    path = tmp_workspace / "test.txt"
    path.write_text("This is a plain text file.\nIt has multiple lines.\na,b,c\n1,2,3")
    return path


# ═══════════════════════════════════════════════════════════════════════
# 1. FILE CONVERTER
# ═══════════════════════════════════════════════════════════════════════


class TestFileConverter:
    def setup_method(self):
        self.mod = _import("file_converter",
                           PROJECT / "actions" / "file_converter.py")

    # helpers exist
    def test_guess_ext_exists(self):
        assert hasattr(self.mod, "_guess_ext")

    def test_output_path_exists(self):
        assert hasattr(self.mod, "_output_path")

    # image -> image
    def test_convert_image_to_jpg(self, with_sample_image):
        result = self.mod.convert_image(str(with_sample_image), ".jpg")
        assert result.endswith(".jpg") and os.path.isfile(result)

    def test_convert_image_to_webp(self, with_sample_image):
        result = self.mod.convert_image(str(with_sample_image), ".webp")
        assert result.endswith(".webp") and os.path.isfile(result)

    def test_convert_image_to_bmp(self, with_sample_image):
        result = self.mod.convert_image(str(with_sample_image), ".bmp")
        assert result.endswith(".bmp") and os.path.isfile(result)

    # image -> PDF
    def test_convert_image_to_pdf(self, with_sample_image):
        result = self.mod.convert_image_to_pdf(str(with_sample_image))
        assert result.endswith(".pdf") and os.path.isfile(result)

    # PDF -> image
    def test_convert_pdf_to_image(self, with_sample_image):
        pdf_path = self.mod.convert_image_to_pdf(str(with_sample_image))
        result = self.mod.convert_pdf_to_image(pdf_path, ".png")
        assert os.path.isfile(result.split("\n")[0])

    # PDF -> DOCX
    def test_convert_pdf_to_docx(self, with_sample_image):
        pdf_path = self.mod.convert_image_to_pdf(str(with_sample_image))
        try:
            result = self.mod.convert_pdf_to_docx(pdf_path)
            assert result.endswith(".docx") and os.path.isfile(result)
        except RuntimeError as e:
            if "PyMuPDF" in str(e):
                pytest.skip("PyMuPDF not installed")

    # TXT -> CSV
    def test_convert_text_to_csv(self, with_sample_txt):
        result = self.mod.convert_text_to_csv(str(with_sample_txt))
        assert result.endswith(".csv") and os.path.isfile(result)

    # CSV -> JSON
    def test_convert_csv_to_json(self, with_sample_csv):
        result = self.mod.convert_csv_to_json(str(with_sample_csv))
        assert result.endswith(".json") and os.path.isfile(result)
        with open(result) as f:
            data = json.load(f)
        assert len(data) == 2 and data[0]["name"] == "Alice"

    # JSON -> CSV
    def test_convert_json_to_csv(self, tmp_workspace):
        p = tmp_workspace / "test.json"
        p.write_text('[{"x":1,"y":2},{"x":3,"y":4}]')
        result = self.mod.convert_json_to_csv(str(p))
        assert result.endswith(".csv") and os.path.isfile(result)

    def test_convert_json_to_csv_empty(self, tmp_workspace):
        p = tmp_workspace / "empty.json"
        p.write_text("[]")
        with pytest.raises(RuntimeError, match="Empty JSON array"):
            self.mod.convert_json_to_csv(str(p))

    # XML -> JSON
    def test_convert_xml_to_json(self, tmp_workspace):
        p = tmp_workspace / "test.xml"
        p.write_text("<root><item><id>1</id><name>Test</name></item></root>")
        result = self.mod.convert_xml_to_json(str(p))
        assert result.endswith(".json") and os.path.isfile(result)

    # MD -> HTML
    def test_convert_md_to_html(self, with_sample_md):
        try:
            result = self.mod.convert_md_to_html(str(with_sample_md))
            assert result.endswith(".html") and os.path.isfile(result)
            assert "h1" in Path(result).read_text()
        except RuntimeError as e:
            if "markdown" in str(e):
                pytest.skip("markdown library not installed")

    # HTML -> PDF
    def test_convert_html_to_pdf(self, tmp_workspace):
        p = tmp_workspace / "test.html"
        p.write_text("<html><body><p>Hello</p></body></html>")
        try:
            result = self.mod.convert_html_to_pdf(str(p))
            assert result.endswith(".pdf") and os.path.isfile(result)
        except RuntimeError as e:
            if "weasyprint" in str(e):
                pytest.skip("weasyprint not installed")

    # DOCX -> PDF
    def test_convert_docx_to_pdf(self, tmp_workspace):
        p = tmp_workspace / "test.docx"
        try:
            from docx import Document
            Document().save(str(p))
        except ImportError:
            pytest.skip("python-docx not installed")
        try:
            result = self.mod.convert_docx_to_pdf(str(p))
            assert result.endswith(".pdf")
        except RuntimeError as e:
            if "weasyprint" in str(e):
                pytest.skip("weasyprint not installed")

    # XLSX -> PDF
    def test_convert_xlsx_to_pdf(self, tmp_workspace):
        p = tmp_workspace / "test.xlsx"
        try:
            from openpyxl import Workbook
            wb = Workbook(); wb.active["A1"] = "Hello"; wb.save(str(p))
        except ImportError:
            pytest.skip("openpyxl not installed")
        try:
            result = self.mod.convert_xlsx_to_pdf(str(p))
            assert result.endswith(".pdf")
        except RuntimeError as e:
            if "weasyprint" in str(e):
                pytest.skip("weasyprint not installed")

    # PPTX -> PDF
    def test_convert_pptx_to_pdf(self, tmp_workspace):
        p = tmp_workspace / "test.pptx"
        try:
            from pptx import Presentation
            prs = Presentation()
            prs.slides.add_slide(prs.slide_layouts[6])
            prs.save(str(p))
        except ImportError:
            pytest.skip("python-pptx not installed")
        try:
            result = self.mod.convert_pptx_to_pdf(str(p))
            assert result.endswith(".pdf")
        except RuntimeError as e:
            if "weasyprint" in str(e):
                pytest.skip("weasyprint not installed")

    # Audio
    def test_convert_audio(self, tmp_workspace):
        wav = tmp_workspace / "test.wav"
        try:
            from pydub import AudioSegment
            AudioSegment.silent(duration=500).export(str(wav), format="wav")
        except ImportError:
            pytest.skip("pydub not installed")
        try:
            result = self.mod.convert_audio(str(wav), ".mp3")
            assert result.endswith(".mp3") and os.path.isfile(result)
        except RuntimeError as e:
            if "pydub" in str(e):
                pytest.skip("pydub not available")

    # OCR
    def test_convert_image_ocr(self, with_sample_image):
        try:
            result = self.mod.convert_image_ocr(str(with_sample_image))
            assert result.endswith(".txt") and os.path.isfile(result)
        except RuntimeError as e:
            if "pytesseract" in str(e):
                pytest.skip("pytesseract not installed")

    # convert_file entry point
    def test_convert_file_missing_source(self):
        r = self.mod.convert_file({})
        assert "No source file" in r or "not specified" in r

    def test_convert_file_not_found(self):
        r = self.mod.convert_file({"source_path": "/nonexistent/foo.png"})
        assert "not found" in r or "does not exist" in r.lower()

    def test_convert_file_image_to_jpg(self, with_sample_image):
        r = self.mod.convert_file({"source_path": str(with_sample_image), "target_format": "jpg"})
        assert "✅" in r or "Converted" in r


# ═══════════════════════════════════════════════════════════════════════
# 2. UNIT CONVERTER
# ═══════════════════════════════════════════════════════════════════════


class TestUnitConverter:
    def setup_method(self):
        self.mod = _import("unit_converter",
                           PROJECT / "actions" / "unit_converter.py")

    def _conv(self, value, from_unit, to_unit):
        return self.mod.convert_units({"value": value, "from": from_unit, "to": to_unit})

    # Length
    def test_miles_to_km(self):
        r = self._conv(26.2, "miles", "km")
        m = re.search(r"= ([\d.]+)", r)
        assert m and abs(float(m.group(1)) - 42.16) < 0.1

    def test_km_to_miles(self):
        r = self._conv(10, "km", "mi")
        assert "mi" in r

    def test_meters_to_feet(self):
        r = self._conv(1, "m", "ft")
        assert "ft" in r or "feet" in r

    # Weight
    def test_kg_to_lb(self):
        r = self._conv(10, "kg", "lb")
        assert "lb" in r

    def test_lb_to_kg(self):
        r = self._conv(22.05, "lb", "kg")
        assert "kg" in r

    # Temperature
    def test_f_to_c(self):
        r = self._conv(100, "f", "c")
        m = re.search(r"= ([\d.]+)", r)
        assert m and abs(float(m.group(1)) - 37.78) < 0.1

    def test_c_to_f(self):
        r = self._conv(0, "c", "f")
        assert "32" in r

    def test_c_to_k(self):
        r = self._conv(0, "c", "k")
        assert "273" in r

    def test_k_to_c(self):
        r = self._conv(273.15, "k", "c")
        assert "0" in r or "0.0" in r

    # Volume
    def test_liters_to_gallons(self):
        r = self._conv(5, "l", "gal")
        assert "gal" in r

    # Speed
    def test_kmh_to_mph(self):
        r = self._conv(100, "km/h", "mph")
        assert "mph" in r or "mi/h" in r

    # Data
    def test_mb_to_gb(self):
        r = self._conv(1024, "mb", "gb")
        assert "gb" in r or "GB" in r

    # Error cases
    def test_invalid_number(self):
        r = self.mod.convert_units({"value": "abc", "from": "km", "to": "mi"})
        assert "Invalid" in r or "not a valid" in r.lower()

    def test_unsupported_conversion(self):
        r = self._conv(10, "km", "celsius")
        assert "Unsupported" in r

    def test_missing_value(self):
        r = self.mod.convert_units({"from": "km", "to": "mi"})
        assert "specify" in r or "value" in r


# ═══════════════════════════════════════════════════════════════════════
# 3. RANDOM NUMBER
# ═══════════════════════════════════════════════════════════════════════


class TestRandomNumber:
    def setup_method(self):
        self.mod = _import("random_number",
                           PROJECT / "actions" / "random_number.py")

    def test_default_range(self):
        for _ in range(50):
            r = self.mod.random_number({})
            m = re.search(r"Random number.*?(\d+)", r)
            assert m, f"No number in: {r}"
            assert 1 <= int(m.group(1)) <= 100

    def test_custom_range(self):
        for _ in range(20):
            r = self.mod.random_number({"min": 50, "max": 60})
            m = re.search(r"(\d+)", r)
            assert m and 50 <= int(m.group(1)) <= 60

    def test_dice(self):
        for _ in range(50):
            r = self.mod.random_number({"action": "dice"})
            m = re.search(r"(\d+)", r)
            assert m and 1 <= int(m.group(1)) <= 6

    def test_coin(self):
        results = set()
        for _ in range(100):
            results.add(self.mod.random_number({"mode": "coin"}).strip())
        assert any("heads" in x.lower() for x in results)
        assert any("tails" in x.lower() for x in results)


# ═══════════════════════════════════════════════════════════════════════
# 4. SYSTEM INFO
# ═══════════════════════════════════════════════════════════════════════


class TestSystemInfo:
    def setup_method(self):
        self.mod = _import("system_info",
                           PROJECT / "actions" / "system_info.py")

    def test_os(self):
        r = self.mod.system_info({"query": "os"})
        assert "Operating System" in r or "OS" in r

    def test_hostname(self):
        r = self.mod.system_info({"query": "hostname"})
        assert "Hostname" in r or "hostname" in r

    def test_cpu(self):
        r = self.mod.system_info({"query": "cpu"})
        assert "Processor" in r or "CPU" in r

    def test_all(self):
        r = self.mod.system_info({"query": "all"})
        assert len(r) > 50

    def test_ram(self):
        r = self.mod.system_info({"query": "ram"})
        assert len(r) > 5  # psutil may be missing but should not crash


# ═══════════════════════════════════════════════════════════════════════
# 5. TODO DISPLAY (parser tests)
# ═══════════════════════════════════════════════════════════════════════


class TestTodoDisplay:
    def setup_method(self):
        self.mod = _import("todo_display",
                           PROJECT / "actions" / "todo_display.py")

    YEAR = "2026"

    def _dt(self, text):
        return self.mod.parse_datetime(text)

    def _pt(self, text):
        return self.mod.parse_task_text(text)

    # parse_datetime
    def test_parse_french_aout(self):
        assert self._dt("21 aout") == f"{self.YEAR}-08-21"

    def test_parse_french_août(self):
        assert self._dt("15 août") == f"{self.YEAR}-08-15"

    def test_parse_french_mars(self):
        assert self._dt("10 mars") == f"{self.YEAR}-03-10"

    def test_parse_english_august(self):
        assert self._dt("august 21") == f"{self.YEAR}-08-21"

    def test_parse_tomorrow(self):
        assert self._dt("tomorrow") != ""

    def test_parse_demain(self):
        assert self._dt("demain") != ""

    def test_parse_in_3_days(self):
        assert self._dt("in 3 days") != ""

    def test_parse_next_week(self):
        assert self._dt("next week") != ""

    def test_parse_friday(self):
        assert self._dt("friday") != ""

    def test_parse_blank(self):
        assert self._dt("") == ""

    def test_parse_gibberish(self):
        assert self._dt("xyzzy flurbo") == ""

    # parse_task_text
    def test_task_title_only(self):
        r = self._pt("buy milk")
        assert r["title"] == "buy milk" and r["due"] == "" and r["priority"] == "normal"

    def test_task_with_due(self):
        r = self._pt("buy milk due tomorrow")
        assert r["title"] == "buy milk" and r["due"] != ""

    def test_task_with_priority(self):
        r = self._pt("write report high priority")
        assert r["title"] == "write report" and r["priority"] == "high"

    def test_task_with_due_and_priority(self):
        r = self._pt("buy milk due tomorrow high priority")
        assert r["title"] == "buy milk" and r["due"] != "" and r["priority"] == "high"

    def test_task_critical(self):
        r = self._pt("fix server critical")
        assert r["priority"] == "critical"

    def test_task_low(self):
        r = self._pt("organize desk low priority")
        assert r["priority"] == "low"

    # show_todo_panel with mock player
    def test_show_todo_panel_with_mock(self, mock_player, isolated_tasks):
        r = self.mod.show_todo_panel(player=mock_player)
        assert r is not None and mock_player.last_url is not None

    def test_show_todo_panel_no_player(self, isolated_tasks):
        r = self.mod.show_todo_panel()
        assert "task" in r.lower() or "No tasks" in r


# ═══════════════════════════════════════════════════════════════════════
# 6. TASK MANAGER
# ═══════════════════════════════════════════════════════════════════════


class TestTaskManager:
    def setup_method(self):
        self.mod = _import("task_manager",
                           PROJECT / "actions" / "task_manager.py")

    def test_crud_cycle(self, isolated_tasks):
        r = self.mod.add_task("crud test")
        assert "added" in r.lower()
        lst = self.mod.list_tasks()
        assert "crud test" in lst
        m = re.search(r"\[(\w+-\d+)\]\s+crud test", lst)
        assert m, f"Task ID not found in:\n{lst}"
        tid = m.group(1)
        r2 = self.mod.complete_task(tid)
        assert "completed" in r2.lower()
        r3 = self.mod.delete_task(tid)
        assert "deleted" in r3.lower()

    def test_add_with_priority_and_due(self, isolated_tasks):
        self.mod.add_task("urgent", priority="high", due="2026-12-31")
        lst = self.mod.list_tasks()
        assert "urgent" in lst and "high" in lst

    def test_complete_not_found(self, isolated_tasks):
        r = self.mod.complete_task("nonexistent-id")
        assert "not found" in r.lower()

    def test_delete_not_found(self, isolated_tasks):
        r = self.mod.delete_task("nonexistent-id")
        assert "not found" in r.lower()

    def test_list_empty(self, isolated_tasks):
        assert "No tasks" in self.mod.list_tasks()

    def test_task_manager_dispatch(self, isolated_tasks):
        r = self.mod.task_manager({"action": "add", "title": "dispatch test"})
        assert "added" in r.lower()

    # Budget
    def test_add_income(self, isolated_tasks):
        r = self.mod.add_transaction("Salary", 5000, ttype="income")
        assert "added" in r.lower()

    def test_add_expense(self, isolated_tasks):
        r = self.mod.add_transaction("Groceries", 150.50, category="food")
        assert "added" in r.lower()

    def test_budget_summary(self, isolated_tasks):
        self.mod.add_transaction("Salary", 5000, ttype="income")
        self.mod.add_transaction("Rent", 1200)
        s = self.mod.budget_summary()
        assert "income" in s.lower() or "expense" in s.lower()

    def test_budget_manager_dispatch(self, isolated_tasks):
        r = self.mod.budget_manager({"action": "add", "description": "Freelance", "amount": 1000, "ttype": "income"})
        assert "added" in r.lower()


# ═══════════════════════════════════════════════════════════════════════
# 7. INTENT ROUTER
# ═══════════════════════════════════════════════════════════════════════


class TestIntentRouter:
    def setup_method(self):
        self.mod = _import("intent_router",
                           PROJECT / "actions" / "intent_router.py")
        self.router = self.mod.get_router()

    def _route(self, text):
        return self.router.route(text)

    def assert_matched(self, text, expected_name=None, expected_params=None):
        result = self._route(text)
        assert result.matched, f"Expected match for: {text!r}"
        if expected_name:
            assert result.intent_name == expected_name, (
                f"Expected intent={expected_name!r}, got {result.intent_name!r} for {text!r}"
            )
        if expected_params:
            for k, v in expected_params.items():
                assert result.handler_params.get(k) == v, (
                    f"Expected param {k}={v!r}, got {result.handler_params.get(k)!r}"
                )
        return result

    def assert_not_matched(self, text):
        result = self._route(text)
        assert not result.matched, f"Should NOT match: {text!r}, got {result.intent_name}"

    # App launcher
    def test_open_chrome(self):
        self.assert_matched("open chrome", "open_app", {"app_name": "chrome"})

    def test_open_google_chrome(self):
        self.assert_matched("open google chrome", "open_app", {"app_name": "google chrome"})

    def test_launch_firefox(self):
        self.assert_matched("launch firefox", "open_app")

    def test_start_vlc(self):
        self.assert_matched("start vlc", "open_app")

    # Email
    def test_read_emails(self):
        self.assert_matched("show my latest emails", "read_emails")

    def test_check_mail(self):
        self.assert_matched("check my email", "read_emails")

    # Weather
    def test_weather_tunis(self):
        self.assert_matched("weather in Tunis", "weather_report", {"city": "Tunis"})

    def test_weather_paris(self):
        self.assert_matched("what is the weather in Paris", "weather_report", {"city": "Paris"})

    # Time / Date
    def test_current_time(self):
        self.assert_matched("what time is it", "get_datetime")

    def test_current_date(self):
        self.assert_matched("what is the date today", "get_datetime")

    # Timer
    def test_set_timer(self):
        self.assert_matched("set a timer for 10 minutes", "set_timer", {"minutes": 10})

    def test_set_timer_5m(self):
        self.assert_matched("set timer 5 minutes", "set_timer", {"minutes": 5})

    def test_alarm(self):
        self.assert_matched("set an alarm for 7 am", "set_timer")

    # Convert units
    def test_convert_miles_to_km(self):
        self.assert_matched("how many km is 26.2 miles", "convert_units")

    def test_convert_f_to_c(self):
        self.assert_matched("convert 100 f to celsius", "convert_units")

    def test_convert_kg_to_lb(self):
        self.assert_matched("what is 10 kg in pounds", "convert_units")

    # Random number
    def test_random_number(self):
        self.assert_matched("generate a random number", "random_number")

    def test_random_number_between(self):
        self.assert_matched("random number between 1 and 100", "random_number")

    def test_roll_dice(self):
        self.assert_matched("roll a dice", "random_number")

    def test_flip_coin(self):
        self.assert_matched("flip a coin", "random_number")

    # System info
    def test_os_info(self):
        self.assert_matched("what is my operating system", "system_info", {"query": "os"})

    def test_cpu_info(self):
        self.assert_matched("what cpu do i have", "system_info", {"query": "cpu"})

    def test_hostname(self):
        self.assert_matched("show hostname", "system_info", {"query": "hostname"})

    def test_ram_info(self):
        self.assert_matched("how much ram do i have", "system_info", {"query": "ram"})

    def test_system_info_all(self):
        self.assert_matched("system info", "system_info", {"query": "all"})

    # Filesystem
    def test_disk_usage(self):
        self.assert_matched("disk usage", "filesystem_query", {"action": "disk_usage"})

    def test_largest_files(self):
        self.assert_matched("largest files in home", "filesystem_query", {"action": "largest"})

    # Todo display
    def test_show_tasks(self):
        self.assert_matched("show my tasks", "todo_display")

    def test_display_todo(self):
        self.assert_matched("display todo list", "todo_display")

    def test_open_tasks(self):
        self.assert_matched("open my tasks", "todo_display")

    # Task add
    def test_add_task_simple(self):
        r = self.assert_matched("add task buy milk", "task_add")
        assert r.handler_params.get("title") == "buy milk"

    def test_add_task_with_due(self):
        r = self.assert_matched("add task buy milk due tomorrow", "task_add")
        assert "title" in r.handler_params

    def test_add_task_with_priority(self):
        r = self.assert_matched("add task write report high priority", "task_add")
        assert r.handler_params.get("priority") == "high"

    # News
    def test_news_tunisia(self):
        self.assert_matched("what the latest tunisia news", "news", {"topic": "tunisia"})

    def test_news_tech(self):
        self.assert_matched("tech news today", "news", {"topic": "tech"})

    def test_news_general(self):
        self.assert_matched("latest news", "news")

    # Real Tutor
    def test_realtutor(self):
        self.assert_matched("open the realtutor", "realtime_tutor")

    def test_real_tutor(self):
        self.assert_matched("real tutor", "realtime_tutor")

    # Web search
    def test_search_web(self):
        self.assert_matched("search the web for python", "web_search")

    def test_google_search(self):
        self.assert_matched("google ai news", "web_search")

    # No-match
    def test_no_match_gibberish(self):
        self.assert_not_matched("xyzzy flurbo garblex")

    def test_no_match_hello(self):
        self.assert_not_matched("hello")

    # Singleton
    def test_router_singleton(self):
        assert self.mod.get_router() is self.mod.get_router()


# ═══════════════════════════════════════════════════════════════════════
# 8. CRASH SAFETY (all modules with bad input)
# ═══════════════════════════════════════════════════════════════════════


class TestCrashSafety:
    def test_file_converter_none(self):
        mod = _import("file_converter", PROJECT / "actions" / "file_converter.py")
        assert isinstance(mod.convert_file(None), str)

    def test_unit_converter_none(self):
        mod = _import("unit_converter", PROJECT / "actions" / "unit_converter.py")
        assert isinstance(mod.convert_units(None), str)

    def test_random_number_none(self):
        mod = _import("random_number", PROJECT / "actions" / "random_number.py")
        assert isinstance(mod.random_number(None), str)

    def test_system_info_none(self):
        mod = _import("system_info", PROJECT / "actions" / "system_info.py")
        assert isinstance(mod.system_info(None), str)

    def test_todo_display_none(self):
        mod = _import("todo_display", PROJECT / "actions" / "todo_display.py")
        assert isinstance(mod.show_todo_panel(None), str)

    def test_task_manager_none(self):
        mod = _import("task_manager", PROJECT / "actions" / "task_manager.py")
        assert isinstance(mod.task_manager(None), str)

    def test_email_reader_none(self):
        mod = _import("email_reader", PROJECT / "actions" / "email_reader.py")
        assert isinstance(mod.read_emails(None), str)

    def test_empty_route(self):
        mod = _import("intent_router", PROJECT / "actions" / "intent_router.py")
        assert not mod.route("").matched

    def test_long_route(self):
        mod = _import("intent_router", PROJECT / "actions" / "intent_router.py")
        assert not mod.route("a" * 10000).matched

    def test_noise_route(self):
        mod = _import("intent_router", PROJECT / "actions" / "intent_router.py")
        assert not mod.route("jarvis jarvis jarvis hello").matched


# ═══════════════════════════════════════════════════════════════════════
# 9. EMAIL READER (credential fallback)
# ═══════════════════════════════════════════════════════════════════════


class TestEmailReader:
    def setup_method(self):
        self.mod = _import("email_reader",
                           PROJECT / "actions" / "email_reader.py")

    def test_no_creds_graceful(self):
        r = self.mod.read_emails({"hours": 1, "limit": 5})
        assert isinstance(r, str) and len(r) > 0

    def test_empty_params(self):
        r = self.mod.read_emails({})
        assert isinstance(r, str)
